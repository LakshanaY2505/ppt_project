from pptx import Presentation            # To create a Powerpoint Presentation

from pptx.util import Inches,Pt          # To specify positions/sizes  
from pptx.dml.color import RGBColor      # To set colours using RGB values 
from pptx.enum.text import PP_ALIGN      # To set text allignment 

from PIL import Image                    # Using Pillow to handle images 
import requests                          # Using Requests library to fetch files from internet 
from io import BytesIO                   # Loads images from requests

from pptx.enum.shapes import MSO_SHAPE   # Contains collection of built-in PowerPoint shapes
from pptx.enum.text import MSO_AUTO_SIZE

# ----------------------
# PPT Creation + Slides 
# ----------------------
# Creating the Powerpoint Presentation
ppt = Presentation()
slide_width, slide_height = ppt.slide_width, ppt.slide_height

#To retrieve the slides (0:title slide, 4:comparison, 5:title only, 6:blank)
layouts = [6,5,6,6,6,4,5,5,6,6,6] 
slides = [ppt.slides.add_slide(ppt.slide_layouts[l]) for l in layouts]

#Adding slides to the presentation 
slide1, slide_toc, slide2, slide3, slide4, slide5, slide6, slide_wordcloud, slide8, slide9, slide_references = slides

# -----------------------------------------------
# HELPER FUNCTIONS : used for reccuring elements 
# -----------------------------------------------
# Convert a measurement from EMUs (English Metric Units) to inches
def emu_to_inches(emu_value):
    return emu_value / 914400  # 1 inch = 914400 EMUs

# Creates rectangle 
def add_rectangle(slide, left, top, width, height, fill="white", outline="white", rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE # allows rounded edges
    # For positioning of rectangle
    rect = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height)) 
    
    fill_style = rect.fill
    fill_style.solid()
    fill_style.fore_color.rgb = COLORS[fill] # fill colour
    rect.line.color.rgb = COLORS[outline]    # outline colour
    return rect

# Creates circle
def add_circle(slide, left, top, size, fill="white", outline="white"):
    # For positioning of circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    # Allows colour to be taken from dictionary or written as rgb colour
    # Fill
    if isinstance(fill, str):      # if a string, get color from COLORS(dict)
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS[fill]
    elif isinstance(fill, RGBColor):  # if already RGBColor, use directly
        circle.fill.solid()
        circle.fill.fore_color.rgb = fill
    # Outline
    if isinstance(outline, str):
        circle.line.color.rgb = COLORS[outline]
    elif isinstance(outline, RGBColor):
        circle.line.color.rgb = outline

    return circle

# Creates textbox 
def add_textbox(slide, left, top, width, height, text, font="bodyFont", color="black", align=PP_ALIGN.CENTER):
    if isinstance(text, str):
        text = [text]  # convert single string to list
    
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)) # for positioning of textbox
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True # enables text wrapping for no overflowing of text
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # autofixes textbox size to fit text
    # if text is within array
    for i, t in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = align

        run = p.add_run()
        run.text = t
        apply_font_style(run, font, color) # apply font style and color
    
    return textbox

# Applying text styling
def apply_font_style(run, font="bodyFont", color="black"):
    style = FONTS["bodyFont"]
    # For flexibilty within text size: 
    if isinstance(font, str):  # allows to give specific font number 
        style = FONTS.get(font, style)
        run.font.name = style["name"]
        run.font.size = style["size"]
        run.font.bold = style["bold"]
        run.font.italic = style["italic"]

    elif isinstance(font, dict): # allows to use the values in the dictionary
        run.font.name = font.get("name", style["name"])
        run.font.size = font.get("size", style["size"])
        run.font.bold = font.get("bold", style["bold"])
        run.font.italic = font.get("italic", style["italic"])

    # Allows string (key) or RGBColor directly
    if isinstance(color, str):
        run.font.color.rgb = COLORS[color]
    else: 
        run.font.color.rgb = color

# Add backgorund image
def add_bg(slide, color=None, image_path=None):
    bg = slide.background.fill
    if image_path: # can add picture
        picture = slide.shapes.add_picture(image_path, 0, 0, ppt.slide_width, ppt.slide_height)
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)
    elif color: # or can make bg a solid colour
        bg.solid()
        bg.fore_color.rgb = COLORS[color]

# -------------------------------
# DICTIONARIES : COLOURS, FONTS  
# -------------------------------
COLORS = {
    "red" : RGBColor(192,0,0),
    "darkRed" : RGBColor(139,0,0),
    "orange" : RGBColor(228, 108, 10),
    "white" : RGBColor(255,255,255),
    "black" : RGBColor(0,0,0),
    "fireBrick" : RGBColor(178, 34, 34),
    "crimson" : RGBColor(90, 3, 10),
    "gold" : RGBColor(252,192,0),
    "flameOrange" : RGBColor(233,113,50),
    "burntOrange" : RGBColor(188,57,8),
    "darkGray" : RGBColor(51,51,51),
    "pumpkinOrange" : RGBColor(230,100,20),
    "tangerine" : RGBColor(232,93,4),
    "oxblood" : RGBColor(56,2,6)
}

FONTS = {
    "titleFont": {
        "name": "Constantia",
        "size": Pt(40),
        "bold": True,
        "italic": False
    },
    "subtitleFont": {
        "name": "Aptos",
        "size": Pt(28),
        "bold": False,
        "italic": True
    },
    "bodyFont": {
        "name": "Aptos",
        "size": Pt(12),
        "bold": False,
        "italic": False
    }
}

# ----------------------------
# Slide 1 : Title slide 
# ----------------------------
add_bg(slide1, color="orange") # Background colour  

# Design for the slide 
# 1-Adding nav bar image
slide1.shapes.add_picture("Images/navBar.png", 0, 0, width=slide_width) 

# 2-Adding circles for cloud design
left1 = 2.5
for j in range(2):
    add_circle(slide1, left1, 5, 2, fill="white", outline="white")
    left1 += 3.5
left2 = 0
for i in range(3):
    add_circle(slide1, left2, 4.5, 3, fill="white", outline="white")
    left2 += 3.5

add_rectangle(slide1, 0, 6, 10, 2, fill="white", outline="white") # Rectangle at bottom

# 3-Centered logo
logo_width = Inches(7)
left = (slide_width - logo_width) / 2
top = Inches(0.2)
slide1.shapes.add_picture("Images/gsLogo.png", left, top, width=logo_width)

# 4.1-Search bar shape
top_position = Inches(2.5)
add_rectangle(slide1, 2, 2.5, 6, 0.5, fill="white", outline="white", rounded=True)
# 4.2-Adding textbox within the shape (with names)
add_textbox(slide1, 2, 2.5, 6, 0.5, "By Lakshana, Kanchan, Tamara", font="subtitleFont", color="black", align=PP_ALIGN.CENTER)

# -------------------------------
# Slide 1.2 : Table of contents 
# -------------------------------
add_bg(slide_toc, color="crimson")  # Background colour  

# Adding title with placeholder
toc = slide_toc.shapes.title
toc.text = "TABLE OF CONTENTS"
p = toc.text_frame.paragraphs[0]
run = p.runs[0]
run.font.name = FONTS["titleFont"]["name"]
run.font.size = FONTS["titleFont"]["size"]
run.font.bold = FONTS["titleFont"]["bold"]
run.font.italic = FONTS["titleFont"]["italic"]
run.font.color.rgb = COLORS["white"]

# Array of TOC items
TOC_array = ["Introduction", "Problem", "Lessons", "Timeline", "Collapse", "Legacy", "Growth", "Cloud", "References"]

# Base positions for rows
row_y_positions = [(1.5, 3),(3.5, 5),(5.5, 7)]  #(image_y, text_y)

images = ["Images/cd1.jpg", "Images/cd2.jpg", "Images/cd3.jpg"]

for i, label in enumerate(TOC_array):
    row = i // 3        # 0,1,2
    col = i % 3         # 0,1,2

    img_y, txt_y = row_y_positions[row]
    x = 2 + col * 2.2

    # Adding the images
    slide_toc.shapes.add_picture(images[row], Inches(x), Inches(img_y), Inches(1.5), Inches(1.5))

    # Add the titles under the pictures
    add_textbox(slide_toc, x - 0.2, txt_y, 1.5, 0.5, label, font="bodyFont", color="white", align=PP_ALIGN.CENTER)

# ----------------------------
# Slide 2 : Introduction 
# ----------------------------
add_bg(slide2, image_path='Images/bg.jpg') # Background colour  

# Adding title
add_textbox(
    slide2, left=1, top=0.5, width=8, height=1, text="INTRODUCTION", font="titleFont", color="red", align=PP_ALIGN.LEFT
)

# Main content shape and textbox
add_rectangle(
    slide2, left=1.25, top=2, width=6, height=4, fill="gold", outline="gold", rounded=True
)
slide2_array = [
    "Grooveshark was an online music streaming service that launched in March 2006,long before Spotify and Apple Music became household names. It offered millions of tracks that users could search, stream, and share freely, making it a favorite hangout spot for music lovers around the world.",
    "Grooveshark was founded by three University of Florida students: Sam Tarantino, Josh Greenberg , Andres Barreto",
    "The founders wanted to make music more accessible to everyone. At a time when piracy was rampant and legal options were limited, they dreamed of creating a community-driven platform where people could upload, share, and enjoy music freely. Their bigger goal was to bridge the gap between listeners and artists in a fairer way than traditional record labels. In essence: music without barriers."
]
content_box = add_textbox(
    slide2, 
    left=1.75, top=2, width=5, height=5, 
    text="", 
    font="bodyFont", 
    color="black", 
    align=PP_ALIGN.LEFT
)
tf = content_box.text_frame
tf.word_wrap = True

for point in slide2_array:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS["black"]
    p.font.name = FONTS["bodyFont"]["name"]
    p.space_after = FONTS["bodyFont"]["size"]

# Adding rectangles for the design of slide
add_rectangle(slide2, 0, 0, 0.5, 7.5, fill="red", outline="red")
add_rectangle(slide2, 0.5, 0, 0.5, 7.5, fill="flameOrange", outline="flameOrange")
add_rectangle(slide2, 8, 0, 2, 7.5, fill="red", outline="red")
add_rectangle(slide2, 9, 0, 1, 7.5, fill="flameOrange", outline="flameOrange")

# Adding images
slide2.shapes.add_picture('Images/vinyl.png', Inches(6.5), Inches(1), Inches(3.5), Inches(6))
slide2.shapes.add_picture('Images/tameimpala.jpg', Inches(8), Inches(2), Inches(4), Inches(4))

# ----------------------------
# Slide 3 : Timeline
# ----------------------------
# Fetching and preparing background image from internet
link = "https://wallpapercave.com/wp/wp4465178.png"
response = requests.get(link)
response.raise_for_status()
img = Image.open(BytesIO(response.content)).convert("RGB")
rotated = img.rotate(180)
cropped = rotated.crop((100, 100, rotated.width - 100, rotated.height - 100))
output_stream = BytesIO()
cropped.save(output_stream, format="JPEG")
output_stream.seek(0)

add_bg(slide3, image_path=output_stream) # Background image

# Adding title
add_textbox(
    slide3, left=1, top=0.5, width=8, height=1, text="TIMELINE", font="titleFont", color="black", align=PP_ALIGN.CENTER
)

# Adding timeline line
add_rectangle(
    slide3, left=5, top=2, width=0.05, height=5, 
    fill="black", outline="black"
)

# Array of events 
events = [
    ("2006-2007","Founded in Florida as a peer-to-peer music sharing site."), 
    ("2007–2008","Shifted into a streaming platform with search, playlists, and community features."), 
    ("2009–2010","Gained popularity- Reached 30M+ users but faced first lawsuits from record labels."),
    ("2011-2012","Hit peak with 35M users, 15M+ songs. Legal pressure grew; apps removed from Apple & Google stores."),
    ("2012–2014","Tried artist-upload model, but lawsuits intensified. Courts ruled executives uploaded copyrighted songs."), 
    ("April 2015","Grooveshark shut down in settlement with Universal, Sony, and Warner.")
]
# Array of colours being used for boxes 
year_colors = [
    RGBColor(255, 165, 0),
    RGBColor(255, 140, 0),
    RGBColor(255, 120, 0),
    RGBColor(255, 80, 0),
    RGBColor(255, 50, 0),
    RGBColor(255, 0, 0)
]

y_start = 2 # position for first event
gap = 1     # vertical gap between events
for i, (year, description) in enumerate(events):
    y_pos = y_start + i * gap
    color = year_colors[i]

    # Event circle
    add_circle(
        slide3, 
        left=4.9, top=y_pos, size=0.25,
        fill=color, # it uses the colours mention in array
        outline=color
    )
    
    # Left/right positions
    if i % 2 == 0:
        x_year = 3.8
        x_desc = 0.8
    else:
        x_year = 5.2
        x_desc = 6.5

    # Year textbox
    add_textbox(
        slide3, 
        left=x_year, top=y_pos-0.3, width=2.5, height=0.5, 
        text=year,
        font="bodyFont", color="black", align=PP_ALIGN.LEFT
    )

    # Description textbox with background color
    desc_box = add_textbox(
        slide3, 
        left=x_desc, top=y_pos-0.6, width=3, height=1, 
        text=description,
        font="bodyFont", color="black", align=PP_ALIGN.LEFT
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = color
    

# ----------------------------
# Slide 4 : Growth
# ----------------------------
add_bg(slide4, image_path='Images/bg.jpg') # Background image

# Adding circles
add_circle(slide4, left=0.35, top=0.35, size=3, fill="red", outline="red")

# Adding title
add_textbox(
    slide4,
    left=0.5, top=1.2, width=3, height=1,
    text="GROWTH",
    font="titleFont",color="black",align=PP_ALIGN.LEFT
)

# Adding content rectangle
add_rectangle(
    slide4,
    left=0.35, top=2, width=5.5, height=4,
    fill="black",outline="black",
    rounded=True
)

add_circle(slide4, left=4.2, top=4, size=3, fill="flameOrange", outline="flameOrange")

slide4.shapes.add_picture('Images/red cover.jpg', Inches(6), Inches(0), Inches(4), Inches(7.5)) # Adding image 

# Array of the content
slide4_array = [
    "Grooveshark worked as a peer-to-peer network, letting users upload MP3s for others to listen to. It later transformed into a sleek web-based music streaming site.",
    "Gained traction quickly because it offered free, on-demand music at a time when few services did Features like search, playlists, and community sharing made it stand out.",
    "Expanded with mobile apps, drawing comparisons to Spotify. By 2011, it had around 35 million users worldwide and a catalog of 15 million+ songs.",
    "Users could follow friends, share playlists, and even discover new artists, giving it a “music community” vibe.",
    "Accessible in over 150 countries, making it one of the most widely used streaming platforms before Spotify took over."
]

# Adding textbox for content
slide4_content = add_textbox(
    slide4,
    left=0.5, top=2, width=5, height=5,
    text="",  
    font="bodyFont",color="white",align=PP_ALIGN.LEFT
)

# Using array, adds all the conent with correct styling 
tf = slide4_content.text_frame
tf.word_wrap = True
for point in slide4_array:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = FONTS["bodyFont"]["size"]
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONTS["bodyFont"]["name"]

# --------------------------------------------
# Slide 5 : Problem
# ~check Charts folder to find bar chart code~
# --------------------------------------------
# Adding title into placeholder with styling
Problem = slide5.shapes.title
Problem.text = "PROBLEM"
run = Problem.text_frame.paragraphs[0].runs[0]
run.font.name = FONTS["titleFont"]["name"]
run.font.size = FONTS["titleFont"]["size"]          
run.font.bold = FONTS["titleFont"]["bold"]             
run.font.color.rgb = COLORS["red"]  

add_bg(slide5, image_path='Images/bg.jpg') # Background image

# Headings for the content and bar chart
title1 = slide5.placeholders[1]
title1.text = "LEGAL FALLOUT"
title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title1.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["red"]

title2 = slide5.placeholders[3]
title2.text = "FEES COMPARISION"
title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title2.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["red"]

# Convert placeholder positions to inches
left = slide5.placeholders[2].left / 914400
top = slide5.placeholders[2].top / 914400
width = slide5.placeholders[2].width / 914400
height = slide5.placeholders[2].height / 914400

# Adding rectangle
content_box = add_rectangle(slide5, 
                            left, top, width, height, 
                            fill="burntOrange", outline="burntOrange", 
                            rounded=True)

# Adding content inside rectangle
tf = content_box.text_frame
tf.clear()
tf.margin_left = Inches(0.2)
tf.margin_right = Inches(0.2)
tf.margin_top = Inches(0.15)
tf.margin_bottom = Inches(0.15)
tf.word_wrap = True

# Adding first point 
p = tf.paragraphs[0]
p.text = "Lack of Licensing - The model relied on user-uploaded content, allowing copyrighted music, uploading without permission and operating in a legally gray area compared to Spotify or Apple Music"
p.font.size = Pt(14)
p.font.name = "Aptos"
p.alignment = PP_ALIGN.LEFT
p.space_after = Pt(6)

# Remaining points
points = [
    "Legal Actions - In 2011, Universal Music Group, Sony Music, Entertainment and Warner Music Group sued Grooveshark for copyright infringement, alleging it ignored takedown notices and deleted the metadata to hide the original uploaders",
    "Court Rulings and Penalties - In 2015, a federal judge deemed Grooveshark's actions willful and in bad faith, resulting in potential $700 million in damages and financial hardship for the company"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(14)
    p.font.name = "Aptos"
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)

# Adding the bar chart image 
img_path = "Images/licensingFees.png"
# Convert placeholder position to inches
ph = slide5.placeholders[4]
left = emu_to_inches(ph.left)
top = emu_to_inches(ph.top)
width = emu_to_inches(ph.width)
height = emu_to_inches(ph.height)
# Adding the iamge to specific positions 
pic = slide5.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))

# ----------------------------
# Slide 6 : Collapse
# ----------------------------
# Adding title into placeholder with styling
Collapse = slide6.shapes.title
Collapse.text = "COLLAPSE"
run = Collapse.text_frame.paragraphs[0].runs[0]
run.font.name = FONTS["titleFont"]["name"]
run.font.size = FONTS["titleFont"]["size"]      
run.font.bold = FONTS["titleFont"]["bold"]             
run.font.color.rgb = COLORS["darkGray"]

# Array of the content 
slide6_array = [
    ("Settlement Agreement", "In April 2015, Grooveshark reached a settlement with Universal, Sony and Warner, requiring immediate cessation of operations, removal of copyrighted music and the transfer of intellectual property to record labels"),
    ("Shutdown & Public Apology", "Grooveshark announced a sudden shutdown, causing users to lose access to playlists and uploaded music. Their apology targeted the artists, songwriters, labels, fans and displaced employees"),
    ("Impact on the Industry", "Millions of consumers were harmed by Grooveshark's forced switch to licensed services such as Spotify and Apple Music, as many lost their favorite tracks and the streaming model they were used to.")
]
# Array of the colours 
colours = [RGBColor(245, 210, 80), RGBColor(219, 124, 38), RGBColor(210, 20, 20)]

# Box dimensions
box_width = Inches(3)
box_height = Inches(3)
spacing = Inches(0.5) #This is space between the boxes

total_width = (3 * box_width) + (2 * spacing) #Calculating total width of all boxes + spacing
start_left = (slide_width - total_width) / 2 #Starting from the left position to center the boxes horizontally
top_pos = Inches(2.5) #This is the vertical position of the boxes from the top of the slide

images = ["Images/yellowVinyl.png", "Images/orangeVinyl.png", "Images/redVinyl.png"]

for i, (title, content) in enumerate(slide6_array): #returns the index value and the tuple itself
    left = start_left + i * (box_width + spacing) #calculates the x position, by using the box width and spacing
    top = top_pos

    # Creates colored rectangle (background)
    box = add_rectangle(slide6, emu_to_inches(left), emu_to_inches(top), emu_to_inches(box_width), emu_to_inches(box_height),
                        fill="white", outline="white")
    box.fill.solid()
    box.fill.fore_color.rgb = colours[i]
    box.line.fill.background()

    # Add internal margins to make sure the text is not touching the edges of the boxes
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_top = Pt(20)      
    tf.margin_bottom = Pt(10)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)

    # Title of the boxes
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = 'Helectivia'
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)        

    # Content of the boxes
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(12)
    p2.font.bold = False
    p2.font.name = FONTS["bodyFont"]["name"]
    p2.font.color.rgb = COLORS["white"]
    p2.alignment = PP_ALIGN.CENTER

    # Adding images by using the array 
    img = images[i]
    img_width = Inches(2)
    img_height = Inches(2)
    img_left = left + (box_width - img_width) / 2 #center horizontally relative to box

    if i in [0, 2]:  #if first or third (will appear on the top)
        img_top = top - Inches(1.3)
    else: #the second image
        img_top = top + box_height - Inches(0.6)

    slide6.shapes.add_picture(img, img_left, img_top, img_width, img_height)


# ---------------------------------------------
# Slide 7 : A Cloud of Controversy (WORDCLOUD)
# ~check Charts folder to find Wordcloud code~
# ---------------------------------------------
add_bg(slide_wordcloud, color="crimson") # Background colour

# Adding title
wordcloud_title = slide_wordcloud.shapes.title
wordcloud_title.text = "A CLOUD OF CONTROVERSY"
run = wordcloud_title.text_frame.paragraphs[0].runs[0]
apply_font_style(run, font="titleFont", color="white")

# Wordcloud image
add_picture_path = "Images/wordcloud.png"
slide_wordcloud.shapes.add_picture(
    add_picture_path,
    Inches(0.5),
    Inches(2),
    Inches(9),
    Inches(4.5)
)

# ---------------------------
# Slide 8 : Lessons Learnt
# ---------------------------
add_bg(slide8, image_path='Images/bg.jpg') # Background image

# Adding title 
add_textbox(slide8, 0.5, 0.8, 9, 1.2, "LESSONS LEARNT", font="titleFont", color="red", align=PP_ALIGN.CENTER)

# Adding content 
subtitle_text = "Despite Grooveshark being an innovative music-sharing platform, it faced challenges in three areas:"
subtitle_box = add_textbox(slide8, 0.8, 2.0, 8.4, 1.0, subtitle_text, font=Pt(20), color="black", align=PP_ALIGN.CENTER)
subtitle_box.text_frame.margin_left = Pt(5)
subtitle_box.text_frame.margin_right = Pt(5)

# Array of music cards with content
card_data = [
    {"title": "Legal Risks", "content": "Allowing users to upload and stream music without proper licenses created significant legal risks.", "has_player": True},
    {"title": "Ethical Concerns", "content": "Although the platform claimed to follow DMCA rules, artists were not paid, raising ethical concerns.", "has_player": True},
    {"title": "Sustainability", "content": "The shutdown highlights that tech innovation must align with ethics and legal requirements to be sustainable.", "has_player": True}
]
# Position of the cards
card_width = 2.8
card_height = 2.0
card_spacing = 0.2
cards_start_top = 3.2
total_cards_width = 3 * card_width + 2 * card_spacing
cards_start_left = (emu_to_inches(slide_width) - total_cards_width) / 2

for i, card in enumerate(card_data):
    card_left = cards_start_left + i * (card_width + card_spacing)
    
    # Title rectangle
    title_box = add_rectangle(slide8, card_left, cards_start_top, card_width, 0.5, fill="pumpkinOrange", rounded=True)
    add_textbox(slide8, card_left, cards_start_top, card_width, 0.7, card["title"], font="bodyFont", color="black", align=PP_ALIGN.CENTER)
    
    # Music controls
    if card["has_player"]:
        player_text = "♥  ⏮  ⏯  ⏭  ⊕"
        add_textbox(slide8, card_left + 0.2, cards_start_top + 0.75, card_width - 0.4, 0.5, player_text, font="bodyFont", color="black", align=PP_ALIGN.CENTER)

    # Content rectangle
    content_top = cards_start_top + 1.2
    content_box = add_rectangle(slide8, card_left, content_top, card_width, 1.6, fill="pumpkinOrange", rounded=True)
    add_textbox(slide8, card_left + 0.15, content_top + 0.15, card_width - 0.3, 1.3, card["content"], font="bodyFont", color="black", align=PP_ALIGN.CENTER)

# -------------------------------------------------
# Slide 9 : Legacy - influence on modern platforms 
# -------------------------------------------------
# Adding title + subtitle
add_textbox(slide9, 1, 0.5, 8, 1.5, ["LEGACY"], font="titleFont", color="darkGray", align=PP_ALIGN.CENTER)
add_textbox(slide9, 1, 1.1, 8, 1.5, [ "~Influence on Modern Platforms~"], font="subtitleFont", color="darkGray", align=PP_ALIGN.CENTER)

# Adding the description content 
desc_text = ("Many features like playlist sharing and music discovery are now common on platforms like Spotify and Youtube music. "
             "Its user-upload model influenced SoundCloud’s community-driven platform. "
             "Became a cautionary tale, proving that long-term success requires proper licensing and artist support")
add_textbox(slide9, 0.5, 2, 9, 1.2, desc_text, font=Pt(20), color="darkGray", align=PP_ALIGN.CENTER)

# Adding table
table_data = [
    ["Feature", "Grooveshark", "Spotify", "YouTube Music", "SoundCloud"],
    ["Playlist Sharing", "✓", "✓", "✓", "✓"],
    ["User Uploads", "✓", "✗", "✗", "✓"],
    ["Recommendations", "✓", "✓", "✓", "✓"],
    ["Licensed Content", "✗", "✓", "✓", "✓ (partial)"],
    ["Artist Payout", "✗", "✓", "✓", "✓ (limited)"]
]

rows = len(table_data)
cols = len(table_data[0])
table = slide9.shapes.add_table(rows, cols, Inches(0.5), Inches(4), Inches(9), Inches(3)).table
# Setting the table (styling, data)
for row_idx, row in enumerate(table_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        # Style header row
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["tangerine"]
            paragraph.font.bold = True
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = COLORS["white"]
        else:
            # Data row styling
            if col_idx == 0:  # Feature name column
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["tangerine"]
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS["white"]
                paragraph.font.size = Pt(14)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 160, 76)
                # Color code the symbols
                if "✓" in val:
                    paragraph.font.color.rgb = RGBColor(0, 128, 0)
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(15)
                elif "✗" in val:
                    paragraph.font.color.rgb = RGBColor(220, 20, 60)
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(15)
                else:
                    paragraph.font.color.rgb = COLORS["darkGray"]
                    paragraph.font.size = Pt(11)
        # Cell margins
        cell.margin_left = Pt(5)
        cell.margin_right = Pt(5)
        cell.margin_top = Pt(5)
        cell.margin_bottom = Pt(5)

table.columns[0].width = Inches(2.2) # feature column wider
for i in range(1, cols):
    table.columns[i].width = Inches(1.7)


# ---------------------------
# Slide 10 : References
# ---------------------------
add_bg(slide_references, image_path='Images/bg.jpg') # Background image

# Adding title 
add_textbox(slide_references, 0.5, 0.5, 9, 1.5, "REFERENCES", font="titleFont", color="red", align=PP_ALIGN.CENTER)

# Container rectangle: centered it horizontally, positioned directly at the bottom of the slide
container_left = (emu_to_inches(slide_width) - 5.5) / 2
container_top = 2
container = add_rectangle(slide_references, container_left, container_top, 5.5, 5, fill="fireBrick", outline="fireBrick", rounded=True)

# Adding an image on top of container 
img_path = "Images/musicControls.png"
img_left = container_left + 1
img_top = container_top - 0.5
img_width = 3.5
img_height = 3

slide_references.shapes.add_picture(img_path, Inches(img_left), Inches(img_top), Inches(img_width), Inches(img_height))

# Reference links
links = [
    "https://www.bbc.com/news/technology-32547376",
    "https://www.rollingstone.com/music/music-news/grooveshark-shuts-down-after-eight-years-74101/",
    "https://www.cinchsolution.com/what-happened-to-grooveshark-com/",
    "https://www.wired.com/story/grooveshark-is-dead/"
]

small_rect_height = 0.6
rect_spacing = 0.05
start_top = container_top + 1.7 # this adds the rectangles below the controls

for i, link in enumerate(links):
    rect_top = start_top + i * (small_rect_height + rect_spacing)
    # 4 smaller rectangles
    rect = add_rectangle(slide_references, container_left, rect_top, 5.5, small_rect_height, fill="darkRed", outline= "darkRed", rounded=False)
    
    # Used to make the hyperlink clickable
    rect.click_action.hyperlink.address = link
    add_textbox(slide_references, container_left + 0.8, rect_top, 5.5-0.8, small_rect_height, link, font="bodyFont", color="white", align=PP_ALIGN.LEFT)
    
    # Circle with number
    circle = add_circle(slide_references, container_left + 0.2, rect_top + (small_rect_height - 0.4)/2, 0.4, fill="oxblood", outline="oxblood")
    add_textbox(slide_references, container_left + 0.2, rect_top + (small_rect_height - 0.4)/2, 0.4, 0.4, str(i+1), font="bodyFont", color="white", align=PP_ALIGN.CENTER)

ppt.save('trial2.pptx') 
