from pptx import Presentation            # To create a Powerpoint Presentation

from pptx.util import Inches,Pt          # To specify positions/sizes  
from pptx.dml.color import RGBColor      # To set colours using RGB values 
from pptx.enum.text import PP_ALIGN      # To set text allignment 

from PIL import Image                    # Using Pillow to handle images 
import requests                          # Using Requests library to fetch files from internet 
from io import BytesIO                   # Loads images from requests

from pptx.enum.shapes import MSO_SHAPE   # Contains collection of built-in PowerPoint shapes

# Creating the Powerpoint Presentation
ppt = Presentation() 
slide_width = ppt.slide_width 
slide_height = ppt.slide_height 

#To retrieve the slides (0:title slide, 4:comparison, 5:title only, 6:blank)
slide1_reg= ppt.slide_layouts[6] 
slide_toc_reg = ppt.slide_layouts[5]
slide2_reg = ppt.slide_layouts[6] 
slide3_reg = ppt.slide_layouts[6] 
slide4_reg = ppt.slide_layouts[6] 
slide5_reg = ppt.slide_layouts[4]
slide6_reg = ppt.slide_layouts[5] 
slide_wordcloud_reg = ppt.slide_layouts[5]
slide8_reg = ppt.slide_layouts[6] 
slide9_reg = ppt.slide_layouts[6]
slide_references_reg = ppt.slide_layouts[5]

#Adding slides to the presentation 
slide1 = ppt.slides.add_slide(slide1_reg) 
slide_toc = ppt.slides.add_slide(slide_toc_reg) 
slide2 = ppt.slides.add_slide(slide2_reg)
slide3 = ppt.slides.add_slide(slide3_reg)
slide4 = ppt.slides.add_slide(slide4_reg)
slide5 = ppt.slides.add_slide(slide5_reg)
slide6 = ppt.slides.add_slide(slide6_reg)
slide_wordcloud = ppt.slides.add_slide(slide_wordcloud_reg)
slide8 = ppt.slides.add_slide(slide8_reg)
slide9 = ppt.slides.add_slide(slide9_reg)
slide_references = ppt.slides.add_slide(slide_references_reg)

# ----------------------------
# Slide 1 : Title slide 
# ----------------------------
#Add a background color 
bg = slide1.background
fill= bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(228, 108, 10)
#Adding circles
left1=2.5
for j in range(2):
    circle = slide1.shapes.add_shape(9, Inches(left1), Inches(5), Inches(2), Inches(2))
    fill = circle.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255,255,255)
    outline = circle.line
    outline.color.rgb = RGBColor(255,255,255)
    left1+=3.5

left2 =0
for i in range(3):
    circle = slide1.shapes.add_shape(9, Inches(left2), Inches(4.5), Inches(3), Inches(3))
    fill = circle.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255,255,255)
    outline = circle.line
    outline.color.rgb = RGBColor(255,255,255)
    left2+=3.5

#Adding a rectangle 
rectangle = slide1.shapes.add_shape(1, Inches(0), Inches(6), Inches(10), Inches(2))
fill = rectangle.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255,255,255)
outline = rectangle.line
outline.color.rgb = RGBColor(255,255,255)

#adding logo 
logo = slide1.shapes.add_picture('Images/groovesharklogo.png', Inches(4), Inches(0.5), Inches(2), Inches(2))

#adding a textbox
textbox = slide1.shapes.add_textbox(Inches(3.7), Inches(2.5), Inches(8), Inches(1.5))
tf=textbox.text_frame
p = tf.paragraphs[0]
p.text="Search for music"
run = p.runs[0] 
p.font.size = Pt(25)
p.font.color.rgb = RGBColor(255,255,255)

#creating a search bar
rect = slide1.shapes.add_shape(5, Inches(2), Inches(3.5), Inches(6), Inches(0.5))
fill = rect.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255,255,255)
line = rect.line
line.color.rgb = RGBColor(255, 255, 255)

#Adding a text box for names
textbox=slide1.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
tf=textbox.text_frame
p=tf.paragraphs[0]
p.text="By Lakshana, Kanchan, Tamara"
runs=p.runs[0]
runs.font.color.rgb = RGBColor(0,0,0)
runs.font.size = Pt(20)

# -------------------------------
# Slide 1.2 : Table of contents 
# -------------------------------
# Adding title 
toc = slide_toc.shapes.title 
toc.text = "TABLE OF CONTENTS"

tf = toc.text_frame   # get text frame inside the shape
p = tf.paragraphs[0]  # get first paragraph
run = p.runs[0]       # get first run (the actual text)
p.font.name = 'Constantia'
run.font.color.rgb = RGBColor(255, 255, 255)
bg=slide_toc.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(90, 3, 10)
x=0.1
y=0.1
z=0.1
TOC_array = ["Introduction", "Problem", "Lessons" , "Timeline", "Collapse","Legacy", "Growth","Cloud", "References"]
for i in range(9):
    if (i%3==0): 
        box = slide_toc.shapes.add_picture("Images/cd1.jpg", Inches(2+x*2.2), Inches(1.5), Inches(1.5), Inches(1.5)) 
        x+=1
        text = slide_toc.shapes.add_textbox(Inches(-0.2+x*2.2), Inches(3), Inches(1.5), Inches(0.5))
        tf = text.text_frame
        tf.text = TOC_array[i]
        p= tf.paragraphs[0]
        runs=p.runs[0]
        runs.font.color.rgb = RGBColor(255,255,255) #setting the font color to white
    elif (i%3==1):
        box = slide_toc.shapes.add_picture("Images/cd2.jpg", Inches(2+y*2.2), Inches(3.5), Inches(1.5), Inches(1.5)) 
        y+=1
        text = slide_toc.shapes.add_textbox(Inches(0+y*2.2), Inches(5), Inches(1.5), Inches(0.5))
        tf = text.text_frame
        tf.text = TOC_array[i]
        p= tf.paragraphs[0]
        runs=p.runs[0]
        runs.font.color.rgb = RGBColor(255,255,255) #setting the font color to white
    else:
        box = slide_toc.shapes.add_picture("Images/cd3.jpg", Inches(2+z*2.2), Inches(5.5), Inches(1.5), Inches(1.5))
        z+=1
        text = slide_toc.shapes.add_textbox(Inches(0+z*2.2), Inches(7), Inches(1.5), Inches(0.5))
        tf = text.text_frame
        tf.text = TOC_array[i]
        p= tf.paragraphs[0]
        runs=p.runs[0]
        runs.font.color.rgb = RGBColor(255,255,255) #setting the font color to white


# ----------------------------
# Slide 2 : Introduction 
# ----------------------------
slide2.shapes.add_picture('Images/bg.jpg', 0, 0, slide_width, slide_height) #adding a bg image (starts from the left=0 inches, top = 0 inches, and spans the full width and height of the screen)

# Adding the title of the slide 
textbox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1)) 
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "INTRODUCTION"
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(192,0,0)
p.font.bold = True
p.font.name = 'Constantia'

# Adding the content: 
# 1- Rectangle design
rectangle = slide2.shapes.add_shape(5, Inches(1.25), Inches (2), Inches(6), Inches(4))
fill=rectangle.fill
fill.solid()
fill.fore_color.rgb = RGBColor(252,192,0)
outline = rectangle.line
outline.color.rgb = RGBColor(252,192,0) 
# 2- Textbox
slide2_array=["Grooveshark was an online music streaming service that launched in March 2006,long before Spotify and Apple Music became household names. It offered millions of tracks that users could search, stream, and share freely, making it a favorite hangout spot for music lovers around the world.", "Grooveshark was founded by three University of Florida students: Sam Tarantino, Josh Greenberg , Andres Barreto", "The founders wanted to make music more accessible to everyone. At a time when piracy was rampant and legal options were limited, they dreamed of creating a community-driven platform where people could upload, share, and enjoy music freely. Their bigger goal was to bridge the gap between listeners and artists in a fairer way than traditional record labels. In essence: music without barriers."]
slide2_content =  slide2.shapes.add_textbox(Inches(1.75), Inches(2), Inches(5), Inches(5))
tf = slide2_content.text_frame
tf.word_wrap = True 
for point in slide2_array: 
    p=tf.add_paragraph()  #adds a new paragraph 
    p.text = point  #sets the text of the paragraph to the current point
    p.font.size = Pt(13)  
    p.font.color.rgb = RGBColor(0,0,0)  #sets the font color to black
    p.font.name = 'Aptos'  #sets the font to Aptos
    p.space_after = Pt(12) 
# 3- Adding rectangles to the design 
rectangle1= slide2.shapes.add_shape(1,Inches(0),Inches(0),Inches(0.5),Inches(7.5))
fill1=rectangle1.fill
fill1.solid()
fill1.fore_color.rgb = RGBColor(192,0,0) #dark red color
outline1 = rectangle1.line
outline1.color.rgb = RGBColor(192,0,0) 

rectangle2= slide2.shapes.add_shape(1,Inches(0.5),Inches(0),Inches(0.5),Inches(7.5))
fill2=rectangle2.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(233,113,50) #dark red color
outline1 = rectangle2.line
outline1.color.rgb = RGBColor(233,113,50) 

rectangle3= slide2.shapes.add_shape(1,Inches(8),Inches(0),Inches(2),Inches(7.5))
fill=rectangle3.fill
fill.solid()
fill.fore_color.rgb = RGBColor(192,0,0) #dark red color
outline1 = rectangle3.line
outline1.color.rgb = RGBColor(192,0,0) 

rectangle4= slide2.shapes.add_shape(1,Inches(9),Inches(0),Inches(1),Inches(7.5))
fill=rectangle4.fill
fill.solid()
fill.fore_color.rgb = RGBColor(233,113,50) #dark red color
outline1 = rectangle4.line
outline1.color.rgb = RGBColor(233,113,50) 

# 4- Adding the image
image=slide2.shapes.add_picture('Images/vinyl.png', Inches(6.5), Inches(1), Inches(3.5), Inches(6)) 
image = slide2.shapes.add_picture('Images/tameimpala.jpg', Inches(8),Inches(2), Inches(4), Inches(4))


# ----------------------------
# Slide 3 : Timeline
# ----------------------------
# Fetching bg image from the internet
link  = "https://wallpapercave.com/wp/wp4465178.png"
response = requests.get(link)
response.raise_for_status()

# Step 1: Open image
img = Image.open(BytesIO(response.content))
# Step 2: Convert to RGB
img = img.convert("RGB")
# Step 3: Rotate
rotated = img.rotate(180)
# Step 4: Crop (optional)
cropped = rotated.crop((100, 100, rotated.width - 100, rotated.height - 100))
# Step 5: Save to memory
output_stream = BytesIO()
cropped.save(output_stream, format="JPEG")
output_stream.seek(0)

# Adding title + content to the slide 
slide3.shapes.add_picture(output_stream, 0, 0, width=slide_width, height=slide_height)
textbox = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
tf = textbox.text_frame
p=tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "TIMELINE"
p.font.size= Pt(40)
p.font.name= 'Constantia'
p.font.bold = True
p.font.color.rgb = RGBColor(0,0,0)  
events= [
    ("2006-2007","Founded in Florida as a peer-to-peer music sharing site."), 
    ("2007–2008","Shifted into a streaming platform with search, playlists, and community features."), 
    ("2009–2010","Gained popularity- Reached 30M+ users but faced first lawsuits from record labels."),
    ("2011-2012","Hit peak with 35M users, 15M+ songs. Legal pressure grew; apps removed from Apple & Google stores."),
    ("2012–2014","Tried artist-upload model, but lawsuits intensified. Courts ruled executives uploaded copyrighted songs."), 
    ("April 2015","Grooveshark shut down in settlement with Universal, Sony, and Warner.")
    ]
# To create a line and adjust the color
line = slide3.shapes.add_shape(1, Inches(5), Inches(2), Inches(0.05), Inches(5))  # rectangle as vertical line (left, top, width,height)
# To color the shape
fill = line.fill #accessing the fill property of the line
fill.solid() #setting the fill to solid
fill.fore_color.rgb = RGBColor(0,0,0) #setting the color of the line to black
# To color the border of the line 
outline = line.line
outline.color.rgb = RGBColor(0,0,0) #setting the color of the border to black
# Adding a starting point for the events 
y_start = 2 #starts at 2 inches from the top (where the line starts)
gap = 1 #gap between each event is 0.7 inches 
y_pos=0  #this is used to calculate the y position needed each time an event is added

year_colors = [
    RGBColor(255, 165, 0),  # orange
    RGBColor(255, 140, 0),  
    RGBColor(255, 120, 0),  
    RGBColor(255, 80, 0),
    RGBColor(255, 50, 0),
    RGBColor(255, 0, 0)     # red
]
for i, (year,description) in enumerate(events):  #enumerate is used so that we can get the index of the event along with the item
    y_pos = y_start + i*gap #calculating the new y position 
    color = year_colors[i]
    #adding a circle for each event 
    circle = slide3.shapes.add_shape(9, Inches(4.9), Inches(y_pos), Inches(0.25), Inches(0.2))  # 9 = oval shape
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    #coloring the border of the circle 
    circle_outline = circle.line 
    circle_outline.color.rgb = color
    if i % 2 == 0: #left side for even index, ride side for odd index
        x_year = 3.8 # inches to the left side
        x_desc = 0.8 
    else:
        x_year = 5.2 # inches to the right side
        x_desc=6.5
    year_box = slide3.shapes.add_textbox(Inches(x_year), Inches(y_pos-0.3), Inches(2.5), Inches(0.5)) #creating a text box for the year
    tf = year_box.text_frame #accessing the text frame of the text box
    p = tf.add_paragraph()  #adding a new paragraph to the text frame
    p.text = year   #setting the text of the paragraph to the year
    p.font.size = Pt(14) 
    p.font.bold = True
    p.font.color.rgb = RGBColor(0,0,0)  #setting the font color to white

    desc_box = slide3.shapes.add_textbox(Inches(x_desc), Inches(y_pos - 0.6), Inches(3), Inches(1))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(14) 
    tf.word_wrap = True
    fill = desc_box.fill
    fill.solid()
    fill.fore_color.rgb = color
    

# ----------------------------
# Slide 4 : Growth
# ----------------------------
# Adding background image 
slide4.shapes.add_picture('Images/bg.jpg', 0, 0, slide_width, slide_height)
# Adding shape 
circle = slide4.shapes.add_shape(9, Inches(0.35), Inches(0.35), Inches(3), Inches(3)) 
fill= circle.fill
fill.solid()
fill.fore_color.rgb = RGBColor(192,0,0) 
outline = circle.line
outline.color.rgb = RGBColor(192,0,0) 
# Adding title 
textbox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.5), Inches(1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "GROWTH"
p.alignment = PP_ALIGN.LEFT
run = p.runs[0]
run.font.size = Pt(40)
run.font.color.rgb = RGBColor(0,0,0)
run.font.bold = True
run.font.name = 'Constantia'

# Adding content
rectangle = slide4.shapes.add_shape(5,Inches(0.35), Inches(2), Inches(5.5), Inches(4))
fill = rectangle.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0,0,0) #black
outline = rectangle.line
outline.color.rgb = RGBColor(0,0,0) #black

circle = slide4.shapes.add_shape(9, Inches(4.2), Inches(4), Inches(3), Inches(3)) 
fill= circle.fill
fill.solid()
fill.fore_color.rgb = RGBColor(233,115,50) 
outline = circle.line
outline.color.rgb = RGBColor(233,115,50) 

image = slide4.shapes.add_picture('Images/red cover.jpg', Inches(6), Inches(0), Inches(4), Inches(7.5)) 

slide4_array = ["Grooveshark worked as a peer-to-peer network,letting users upload MP3s for others to listen to. It later transformed into a sleek web-based music streaming site.", "Gained traction quickly because it offered free, on-demand music at a time when few services did Features like search, playlists, and community sharing made it stand out.","Expanded with mobile apps, drawing comparisons to Spotify. By 2011, it had around 35 million users worldwide and a catalog of 15 million+ songs.","Users could follow friends, share playlists, and even discover new artists, giving it a “music community” vibe.","Accessible in over 150 countries, making it one of the most widely used streaming platforms before Spotify took over."]
slide4_content  = slide4.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(5))  #adding a textbox to the slide

for point in slide4_array:
    tf = slide4_content.text_frame
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255,255,255) 
    tf.word_wrap = True 

# ---------------------------------------------
# Slide 5 : Problem
# ~check Charts folder to find bar chart code~
# ---------------------------------------------
# Adding the title
Problem = slide5.shapes.title
Problem.text = "PROBLEM"

run = Problem.text_frame.paragraphs[0].runs[0]

# Set font properties
run.font.name = "Constantia"
run.font.size = Pt(40)          
run.font.bold = True             
run.font.color.rgb = RGBColor(192, 0, 0)  


# Adding the background image and sending it to the back
background_image="Images/bg.jpg"
background=slide5.shapes.add_picture(
    background_image,
    0,0, #Starts from the (0,0) position
    width=ppt.slide_width,
    height=ppt.slide_height
)

# Adding the background image to the bottom of the stack
slide5.shapes._spTree.remove(background._element)
slide5.shapes._spTree.insert(2, background._element)

# Adding the 2 subtitles
title1 = slide5.placeholders[1]
title1.text = "LEGAL FALLOUT"
title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title1.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0)

title2 = slide5.placeholders[3]
title2.text = "FEES COMPARISION"
title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0)

# Creating the rounded shape
text_placeholder = slide5.placeholders[2]  #Left side placeholder position
left = text_placeholder.left
top = text_placeholder.top
width = text_placeholder.width
height = text_placeholder.height

content_box = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    left,
    top,
    width,
    height
)

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(188,57,8)  
content_box.line.color.rgb = RGBColor(188,57,8) #Same color border        

# Adding the text inside the rectangle
text_frame = content_box.text_frame
text_frame.clear()
text_frame.margin_left = Inches(0.2)
text_frame.margin_right = Inches(0.2)
text_frame.margin_top = Inches(0.15)
text_frame.margin_bottom = Inches(0.15)
text_frame.word_wrap = True

# Add bullets
p = text_frame.paragraphs[0]
p.text = "Lack of Licensing - The model relied on user-uploaded content, allowing copyrighted music, uploading without permission and operating in a legally gray area compared to Spotify or Apple Music"
p.font.size = Pt(14)
p.font.name = "aptos"
p.level = 0
p.space_after = Pt(6)
p.alignment=PP_ALIGN.LEFT

points = [
    "Legal Actions - In 2011, Universal Music Group, Sony Music, Entertainment and Warner Music Group sued Grooveshark for copyright infringement, alleging it ignored takedown notices and deleted the metadata to hide the original uploaders",
    "Court Rulings and Penalties - In 2015, a federal judge deemed Grooveshark's actions willful and in bad faith, resulting in potential $700 million in damages and financial hardship for the company"
]

for point in points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(14)
    p.font.name = "aptos"
    p.level = 0 #Top level bullet - no indentation
    p.space_after = Pt(6) #Adds 6 points of vertical spacing after the paragraph
    p.alignment=PP_ALIGN.LEFT

# Adding the image to the right-side placeholder
content_placeholder = slide5.placeholders[4]  
img_path = "Images/licensingFees.png" 

left = content_placeholder.left
top = content_placeholder.top
width = content_placeholder.width
height = content_placeholder.height

slide5.shapes.add_picture(img_path, left, top, width, height)


# ----------------------------
# Slide 6 : Collapse
# ----------------------------
# Adding the title
Collapse = slide6.shapes.title
Collapse.text = "COLLAPSE"

run = Collapse.text_frame.paragraphs[0].runs[0]

# Set font properties
run.font.name = 'Constantia'
run.font.size = Pt(40)      
run.font.bold = True             
run.font.color.rgb = RGBColor(51, 51, 51)

slide6_array = [ #3 tuples, the first element is the title and the second is the content
    ("Settlement Agreement", "In April 2015, Grooveshark reached a settlement with Universal, Sony and Warner, requiring immediate cessation of operations, removal of copyrighted music and the transfer of intellectual property to record labels"),
    ("Shutdown & Public Apology", "Grooveshark announced a sudden shutdown, causing users to lose access to playlists and uploaded music. Their apology targeted the artists, songwriters, labels, fans and displaced employees"),
    ("Impact on the Industry", "Millions of consumers were harmed by Grooveshark's forced switch to licensed services such as Spotify and Apple Music, as many lost their favorite tracks and the streaming model they were used to.")
]

colours = [ # Colors of the boxes
    RGBColor(245, 210, 80), #Yellow
    RGBColor(219, 124, 38), #Orange
    RGBColor(210, 20, 20)   #Red
]

# Box dimensions
box_width = Inches(3)
box_height = Inches(3)
spacing = Inches(0.5)  #This is space between the boxes

# Calculating total width of all boxes + spacing
total_width = (3 * box_width) + (2 * spacing)

# Starting from the left position to center the boxes horizontally
start_left = (slide_width - total_width) / 2

#This is the vertical position of the boxes from the top of the slide
top_pos = Inches(2.5)

# Image paths
images = [
    "Images/yellowVinyl.png",  #Will stick out from the top
    "Images/orangeVinyl.png",  #Will stick out from the bottom
    "Images/redVinyl.png"   #Will stick out from the top
]

for i, (title, content) in enumerate(slide6_array): #returns the index value and the tuple itself
    left = start_left + i * (box_width + spacing) #calculates the x position, by using the box width and spacing
    top = top_pos #y-axis remains the same

    # Create colored rectangle (background)
    box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height) #Adds the rectangle, according to the dimensions
    box.fill.solid()
    box.fill.fore_color.rgb = colours[i] #Choose from the colors array
    box.line.fill.background()  #Removes the border from the boxes

    # Add text inside rectangle
    text_frame = box.text_frame
    text_frame.clear() #Clear any previous text

    # Title of the boxes
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = 'Helectivia'
    p.font.bold=True
    p.font.color.rgb = RGBColor(255, 255, 255) 
    p.alignment = PP_ALIGN.CENTER

    # Content of the boxes
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(12)
    p.font.bold = False
    p.font.name = 'aptos'
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Add internal margins to make sure the text is not touching the edges of the boxes
    text_frame.margin_top = Pt(10)
    text_frame.margin_bottom = Pt(10)
    text_frame.margin_left = Pt(10)
    text_frame.margin_right = Pt(10)

    # Chooses the image from the array
    img = images[i] 
    # Image dimensions
    img_width = Inches(2)
    img_height = Inches(2)
    # Calculate the image position
    img_left = left + (box_width - img_width) / 2  #center horizontally relative to box

    if i == 0 or i == 2:  #if first or third (will appear on the top)
        img_top = top - Inches(1.3)  #adjust accordingly (2 boxes)
    elif i == 1:  #the second image
        img_top = top + box_height - Inches(0.6)  #will appear on the bottom 

    slide6.shapes.add_picture(img, img_left, img_top, img_width, img_height) #adds the image

# -----------------------------------------------
# Slide 7 : A Cloud of Controversy (WORDCLOUD)
# ~check Charts folder to find Wordcloud code~
# -----------------------------------------------
fill = slide_wordcloud.background.fill
fill.solid()  #Solid color for the background
fill.fore_color.rgb = RGBColor(90, 3, 10)


wordcloud_image="Images/wordcloud.png"
slide_wordcloud.shapes.add_picture(
    wordcloud_image,
    Inches(0.5), #left
    Inches(2), #top
    Inches(9), #width
    Inches(4.5) #height
)

wordcloud_title=slide_wordcloud.shapes.title
wordcloud_title.text="A CLOUD OF CONTROVERSY"
wordcloud_title.text_frame.paragraphs[0].runs[0].font.name = 'constantia' #paragraph[0] is the first paragraph of information
wordcloud_title.text_frame.paragraphs[0].runs[0].font.color.rgb=RGBColor(255, 255, 255)


# ---------------------------
# Slide 8 : Lessons Learnt
# ---------------------------
slide8.shapes.add_picture("Images/bg.jpg", Inches(0), Inches(0), slide_width, slide_height)

# Adding title to slide 
title_left = Inches(0.5)
title_top = Inches(0.8)
title_width = Inches(9)
title_height = Inches(1.2)

title_box = slide8.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame

title_frame.clear()

p = title_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER  

run = p.add_run()
run.text = "LESSONS LEARNT"    

font = run.font
font.name = 'Constantia'
font.size = Pt(40)
run.font.bold = True    
font.color.rgb = RGBColor(192, 0, 0)

# Adding content 
subtitle_left = Inches(0.8)
subtitle_top = Inches(2.0)
subtitle_width = Inches(8.4)
subtitle_height = Inches(1.0)

subtitle_box = slide8.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Despite Grooveshark being an innovative music-sharing platform, it faced challenges in three areas:"

subtitle_frame.word_wrap = True
subtitle_frame.margin_left = Pt(5)
subtitle_frame.margin_right = Pt(5)

subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.alignment = PP_ALIGN.CENTER
subtitle_font = subtitle_paragraph.font
subtitle_font.name = 'Arial'
subtitle_font.size = Pt(20)
subtitle_font.color.rgb = RGBColor(0, 0, 0)  # Black text


# Music cards with content
card_data = [
    {
        "title": "Legal Risks",
        "content": "Allowing users to upload and stream music without proper licenses created significant legal risks.",
        "has_player": True
    },
    {
        "title": "Ethical Concerns",
        "content": "Although the platform claimed to follow DMCA rules, artists were not paid, raising ethical concerns.",
        "has_player": True
    },
    {
        "title": "Sustainability",
        "content": "The shutdown highlights that tech innovation must align with ethics and legal requirements to be sustainable.",
        "has_player": True
    }
]

# Position of the cards
card_width = Inches(2.8)
card_height = Inches(2.0)
card_spacing = Inches(0.2)
cards_start_top = Inches(3.2)

total_cards_width = 3 * card_width + 2 * card_spacing
cards_start_left = (slide_width - total_cards_width) / 2

# Titles for cards
for i, data in enumerate(card_data):
    card_left = cards_start_left + i * (card_width + card_spacing)
    
    title_card_height = Inches(0.7)
    title_card_shape = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        card_left, cards_start_top, card_width, title_card_height
    )
    
    # Adding boxes for title cards
    title_fill = title_card_shape.fill
    title_fill.solid()
    title_fill.fore_color.rgb = RGBColor(230, 100, 20)  
    
    title_line = title_card_shape.line # removing border
    title_line.fill.background()
    
    # Adding title text + format
    title_text_frame = title_card_shape.text_frame
    title_text_frame.clear()
    title_text_frame.margin_left = Inches(0.1)
    title_text_frame.margin_right = Inches(0.1)
    title_text_frame.margin_top = Inches(0.15)
    title_text_frame.margin_bottom = Inches(0.1)
    
    p = title_text_frame.paragraphs[0]
    p.text = data["title"]
    p.alignment = PP_ALIGN.CENTER
    
    font = p.font
    font.name = 'Arial'
    font.size = Pt(16)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0) 
    
    # Adding music player controls right below title
    if data["has_player"]:
        # Position for music player controls 
        player_width = Inches(2.4)
        player_height = Inches(0.5)
        player_left = card_left + (card_width - player_width) / 2
        player_top = cards_start_top + title_card_height + Inches(0.05)
        
        # Adding a text box for the controls
        player_box = slide8.shapes.add_textbox(
            player_left, player_top, player_width, player_height
        )
        player_frame = player_box.text_frame
        player_frame.text = "♥  ⏮  ⏯  ⏭  ⊕"  
        
        player_paragraph = player_frame.paragraphs[0]
        player_paragraph.alignment = PP_ALIGN.CENTER
        player_font = player_paragraph.font
        player_font.name = 'Arial'
        player_font.size = Pt(20)
        player_font.color.rgb = RGBColor(0, 0, 0)
    
    # Adding the content card below music controls
    content_card_top = cards_start_top + title_card_height + Inches(0.6)  # Space for music controls
    content_card_height = Inches(1.6)  # Adjusted height
    
    content_card_shape = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        card_left, content_card_top, card_width, content_card_height
    )
    
    # Adding the boxes for the content boxes
    content_fill = content_card_shape.fill
    content_fill.solid()
    content_fill.fore_color.rgb = RGBColor(230, 100, 20)  # Orange background
    
    content_line = content_card_shape.line
    content_line.fill.background()
    
    # Add card content text + format
    card_text_frame = content_card_shape.text_frame
    card_text_frame.clear()  # Clear default text
    card_text_frame.margin_left = Inches(0.15)
    card_text_frame.margin_right = Inches(0.15)
    card_text_frame.margin_top = Inches(0.15)
    card_text_frame.margin_bottom = Inches(0.15)
    card_text_frame.word_wrap = True
    
    p = card_text_frame.paragraphs[0]
    p.text = data["content"]
    p.alignment = PP_ALIGN.CENTER
    
    font = p.font
    font.name = 'Arial'
    font.size = Pt(14)
    font.color.rgb = RGBColor(0,0,0)  

# -------------------------------------------------
# Slide 9 : Legacy - influence on modern platforms 
# -------------------------------------------------
# Adding title + format
title_box = slide9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame

title_frame.clear()

# --- Main Title ---
p1 = title_frame.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
run1 = p1.add_run()
run1.text = "LEGACY"
font1 = run1.font
font1.name = 'Constantia'
font1.size = Pt(40)
font1.bold = True
font1.color.rgb = RGBColor(51, 51, 51)

# --- Subtitle ---
p2 = title_frame.add_paragraph()  # create a new paragraph
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "~Influence on Modern Platforms~"
font2 = run2.font
font2.name = 'Constantia'
font2.size = Pt(28)
font2.bold = True
font2.color.rgb = RGBColor(51, 51, 51)


# Adding content 
desc_box = slide9.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
desc_frame = desc_box.text_frame
desc_frame.text = ("Many features like playlist sharing and music discovery are now common on platforms like Spotify and Youtube music. ​Its user-upload model influenced SoundCloud’s community-driven platform​. Became a cautionary tale, proving that long-term success requires proper licensing and artist support")
desc_frame.word_wrap = True  # enable word wrapping
desc_paragraph = desc_frame.paragraphs[0]
desc_paragraph.alignment = PP_ALIGN.CENTER
desc_paragraph.font.size = Pt(20)
desc_paragraph.font.color.rgb = RGBColor(51, 51, 51)  

# Table array 
table_data = [
    ["Feature", "Grooveshark", "Spotify", "YouTube Music", "SoundCloud"],
    ["Playlist Sharing", "✓", "✓", "✓", "✓"],
    ["User Uploads", "✓", "✗", "✗", "✓"],
    ["Recommendations", "✓", "✓", "✓", "✓"],
    ["Licensed Content", "✗", "✓", "✓", "✓ (partial)"],
    ["Artist Payout", "✗", "✓", "✓", "✓ (limited)"]
]

# Adding table 
rows = len(table_data)
cols = len(table_data[0])
table = slide9.shapes.add_table(rows, cols, Inches(0.5), Inches(4), Inches(9), Inches(3)).table # positioning it under paragraph 

# Setting the table (styling, data)
for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            
            # Style header row
            if row_idx == 0:
                # Header styling
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 93, 4) # Orange
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(15)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Data row styling
                cell.fill.solid()
                if col_idx == 0:  # Feature name column
                    cell.fill.fore_color.rgb = RGBColor(232, 93, 4)  # Light orange
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(14)
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    cell.fill.fore_color.rgb = RGBColor(248, 160, 76)  # Very light orange
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(11)
                    
                    # Color code the symbols
                    if "✓" in cell_data:
                        paragraph.font.color.rgb = RGBColor(0, 128, 0)  # Green
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(15) 
                    elif "✗" in cell_data:
                        paragraph.font.color.rgb = RGBColor(220, 20, 60)  # Crimson red
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(15) 
                    else:
                        paragraph.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
                
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Cell margins
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(5)
            cell.margin_bottom = Pt(5)

table.columns[0].width = Inches(2.2)  # feature column wider
for i in range(1, cols):
    table.columns[i].width = Inches(1.7)       


# ---------------------------
# Slide 10 : References
# ---------------------------
# Adding the background image 
background_image="Images/bg.jpg"
slide_references.shapes.add_picture(background_image,0,0,width=slide_width, height=slide_height)

# Adding title + format
References = slide_references.shapes.add_textbox(
    Inches(0.5),  # left
    Inches(0.5),  # top
    Inches(9),    # width
    Inches(1.5)   # height
)
References.text_frame.text = "REFERENCES"
References.text_frame.paragraphs[0].font.size = Pt(40)
References.text_frame.paragraphs[0].font.color.rgb = RGBColor(192,0,0)
References.text_frame.paragraphs[0].font.name="Constantia"
References.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Adding rounded rectangle shape
container_width = Inches(5.5)
container_height = Inches(5)
container_left = (slide_width - container_width) / 2 # used to center it horizontally
container_top = Inches(2) # to position directly at the bottom of the slide

container = slide_references.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    container_left,
    container_top,
    container_width,
    container_height
)
container.fill.solid()
container.fill.fore_color.rgb = RGBColor(178, 34, 34)  
container.line.color.rgb=RGBColor(178, 34, 34)   


img_path = "Images/musicControls.png"  
img_width = container_width - Inches(2)
img_height = Inches(3)
img_left = container_left + Inches(1)
img_top = container_top - Inches(0.5)

slide_references.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

small_rect_width = container_width
small_rect_height = Inches(0.6)
rect_spacing = Inches(0.05)
start_top = container_top + Inches(1.7)  # this adds the rectangles below the controls

# List of reference links
links = [
    "https://www.bbc.com/news/technology-32547376",
    "https://www.rollingstone.com/music/music-news/grooveshark-shuts-down-after-eight-years-74101/",
    "https://www.cinchsolution.com/what-happened-to-grooveshark-com/",
    "https://www.wired.com/story/grooveshark-is-dead/"
]

for i, link in enumerate(links):
    rect_top = start_top + i * (small_rect_height + rect_spacing)

    # 4 smaller rectangles
    rect = slide_references.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        container_left,
        rect_top,
        small_rect_width,
        small_rect_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(139, 0, 0) 
    rect.line.color.rgb = RGBColor(139, 0, 0) 

    # Used to make the hyperlink clickable
    rect.click_action.hyperlink.address = link

    # Adding the the hyperlink inside the rectangles
    text_frame = rect.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = link
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)  
    p.alignment = PP_ALIGN.LEFT
    text_frame.margin_left = Inches(0.8)  #used to add space for the circles

    # Adding the circle
    circle_radius = Inches(0.4)
    circle = slide_references.shapes.add_shape(
        MSO_SHAPE.OVAL,
        container_left + Inches(0.2),
        rect_top + (small_rect_height - circle_radius) / 2,
        circle_radius,
        circle_radius
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(56, 2, 6)
    circle.line.color.rgb=RGBColor(56, 2, 6)

    # Adding the numbers
    circle_text = circle.text_frame
    circle_text.text = str(i+1)
    circle_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    circle_text.paragraphs[0].font.size = Pt(14)
    circle_text.paragraphs[0].alignment = PP_ALIGN.CENTER


ppt.save('Grooveshark.pptx') 
